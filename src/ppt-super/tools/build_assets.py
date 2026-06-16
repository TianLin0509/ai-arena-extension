# ppt-super/tools/build_assets.py — 为 PPT-SUPER 准备首版 8 套模板资源
# 产物：src/ppt-super/assets/<id>/{template.pptx, blank.png, thumb.png}
#       src/ppt-super/templates.json （8 套索引 + 每套字段定义）
# 跑法：uv run python C:\Users\lintian\AI_debate\ai-arena-extension\src\ppt-super\tools\build_assets.py
import json, re, shutil
from pathlib import Path
from pptx import Presentation

LIB = Path(r"C:\Users\lintian\ppt-assistant\library")
SUPER = Path(r"C:\Users\lintian\AI_debate\ai-arena-extension\src\ppt-super")
ASSETS = SUPER / "assets"
IDS = ["tpl-01-quad-grid", "tpl-02-stage-evolution", "tpl-03-tri-evolution", "tpl-04-asis-tobe",
       "tpl-05-staircase", "tpl-06-radial-hub", "tpl-07-loop-flywheel", "tpl-08-funnel",
       "tpl-09-venn-fusion", "tpl-10-three-pillar", "tpl-11-four-cards", "tpl-12-mirror-compare",
       "tpl-13-pyramid", "tpl-14-timeline-roadmap", "tpl-15-arch-layers", "tpl-16-kpi-dashboard",
       "tpl-17-problem-solution", "tpl-18-core-stacked", "tpl-19-topic-table", "tpl-20-summary-next",
       "tpl-21-scene-spotlight", "tpl-22-hero-split", "tpl-28-product-center", "tpl-35-cover"]
SLOT_FIELDS = ("key", "type", "role", "zh", "bbox", "font_pt", "chars", "single_line", "hint", "sample")


def iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from iter_shapes(sh.shapes)
        else:
            yield sh


def extract_style(sh):
    """从 [key] 形状的 XML 提取预览所需真实样式：水平/垂直对齐 + 字号 + 加粗 + 文字色。
    预览浮层据此精确还原，消除「左上顶对齐黑字」与 PowerPoint 实渲的偏差。"""
    xml = sh._element.xml

    def f(pat, default=None):
        m = re.search(pat, xml)
        return m.group(1) if m else default

    style = {
        "align": f(r'algn="(\w+)"', "l"),       # l / ctr / r / just（默认左）
        "anchor": f(r'anchor="(\w+)"', "t"),     # t / ctr / b（默认顶）
        "bold": 1 if re.search(r'\bb="1"', xml) else 0,
    }
    sz = f(r'sz="(\d+)"')                          # 百分之一磅
    if sz:
        style["sz_pt"] = round(int(sz) / 100, 1)
    clr = f(r'<a:srgbClr val="([0-9A-Fa-f]{6})"')  # 文字色（模板文字框无背景填充，首个即文字色）
    if clr:
        style["color"] = clr
    return style


def render_png(pptx_in, png_out):
    import pythoncom, win32com.client
    pythoncom.CoInitialize()
    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(str(pptx_in), ReadOnly=False, Untitled=False, WithWindow=False)
    try:
        pres.Slides(1).Export(str(png_out), "PNG", 1920, 1080)
    finally:
        pres.Close()


def main():
    templates = []
    for tid in IDS:
        src, dst = LIB / tid, ASSETS / tid
        dst.mkdir(parents=True, exist_ok=True)
        shutil.copy(src / "template.pptx", dst / "template.pptx")
        if (src / "preview.png").exists():
            shutil.copy(src / "preview.png", dst / "thumb.png")
        # 空背景图：清空 [key] 文字（保留模板自带固定结构文字如 As-Is/VS），再 COM 渲染
        prs = Presentation(str(src / "template.pptx"))
        for sh in iter_shapes(prs.slides[0].shapes):
            if sh.name.startswith("[") and sh.name.endswith("]") and sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        r.text = ""
        tmp = dst / "_blank.pptx"
        prs.save(str(tmp))
        try:
            render_png(tmp, dst / "blank.png")
            print(f"  [{tid}] blank.png OK", flush=True)
        except Exception as e:
            print(f"  [{tid}] blank render FAIL: {e}", flush=True)
        finally:
            try:
                tmp.unlink()
            except Exception:
                pass
        desc = json.loads((src / "descriptor.json").read_text(encoding="utf-8"))
        smap = {sh.name: sh for sh in iter_shapes(prs.slides[0].shapes)}
        slots = []
        for s in desc["slots"]:
            slot = {k: s.get(k) for k in SLOT_FIELDS}
            sh = smap.get("[" + (s.get("key") or "") + "]")
            if s.get("type") != "image" and sh is not None and sh.has_text_frame:
                slot["style"] = extract_style(sh)
            slots.append(slot)
        templates.append({
            "id": tid, "name_cn": desc.get("name_cn", tid),
            "when_to_use": desc.get("when_to_use", ""),
            "canvas": desc.get("canvas", [1920, 1080]),
            "thumb": f"ppt-super/assets/{tid}/thumb.png",
            "blank": f"ppt-super/assets/{tid}/blank.png",
            "pptx": f"ppt-super/assets/{tid}/template.pptx",
            "slots": slots,
        })
    out = SUPER / "templates.json"
    out.write_text(json.dumps({"version": 1, "templates": templates}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"done: {len(templates)} templates -> {out}", flush=True)


if __name__ == "__main__":
    main()
