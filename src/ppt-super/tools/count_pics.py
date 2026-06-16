# 数 pptx 第一页的图片(<p:pic>)数量 + 文字 slot 填充情况
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

P = Path(sys.argv[1] if len(sys.argv) > 1 else
         r"C:\Users\lintian\AI_debate\ai-arena-extension\src\ppt-super\tools\_test_icon_out.pptx")


def it(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from it(sh.shapes)
        else:
            yield sh


prs = Presentation(str(P))
pics, txt = 0, 0
for sh in it(prs.slides[0].shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pics += 1
    if sh.name.startswith("[") and sh.name.endswith("]") and sh.has_text_frame and sh.text_frame.text.strip():
        txt += 1
print(f"pics(<p:pic>): {pics}   filled text slots: {txt}")
print("OPEN OK")
