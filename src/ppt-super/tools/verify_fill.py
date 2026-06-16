# 读回 JS 引擎生成的 pptx，确认能被 python-pptx 打开（无损坏）+ 文字填对
import sys
from pathlib import Path
from pptx import Presentation

P = Path(sys.argv[1] if len(sys.argv) > 1 else
         r"C:\Users\lintian\AI_debate\ai-arena-extension\src\ppt-super\tools\_test_out.pptx")


def it(shapes):
    for sh in shapes:
        if sh.shape_type == 6:
            yield from it(sh.shapes)
        else:
            yield sh


prs = Presentation(str(P))
n = 0
empty = []
for sh in it(prs.slides[0].shapes):
    if sh.name.startswith("[") and sh.name.endswith("]") and sh.has_text_frame:
        t = sh.text_frame.text
        print(f"{sh.name:24} {t[:44]!r}")
        n += 1
        if not t.strip():
            empty.append(sh.name)
print(f"\nfilled [key] slots: {n}   empty: {empty}")
print("OPEN OK — python-pptx parsed without error")
