# bake_one.py — 增量烤一个 codex-library/library 模板进 PPT-SUPER，追加/覆盖 templates.json
# 跑法: uv run --with python-pptx --with pywin32 python bake_one.py [模板目录绝对路径 ...]
#   不带参数默认烤 tpl-99-analysis-solution-split
# 核心(pptx/thumb/slots/templates.json)不依赖 COM；blank.png 用 COM 渲染，失败用 thumb 兜底不阻塞。
import json, re, shutil, sys
from pathlib import Path
from pptx import Presentation

SUPER = Path(r"C:\Users\lintian\AI_debate\ai-arena-extension\src\ppt-super")
ASSETS = SUPER / "assets"
SLOT_FIELDS = ("key", "type", "role", "zh", "bbox", "font_pt", "chars", "single_line", "hint", "sample")
DEFAULT_SRC = r"C:\Users\lintian\ppt-assistant\codex-library\tpl-99-analysis-solution-split"


def iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from iter_shapes(sh.shapes)
        else:
            yield sh


def extract_style(sh):
    xml = sh._element.xml

    def f(pat, d=None):
        m = re.search(pat, xml)
        return m.group(1) if m else d

    style = {
        "align": f(r'algn="(\w+)"', "l"),
        "anchor": f(r'anchor="(\w+)"', "t"),
        "bold": 1 if re.search(r'\bb="1"', xml) else 0,
    }
    sz = f(r'sz="(\d+)"')
    if sz:
        style["sz_pt"] = round(int(sz) / 100, 1)
    clr = f(r'<a:srgbClr val="([0-9A-Fa-f]{6})"')
    if clr:
        style["color"] = clr
    return style


def render_blank_png(pptx_in, png_out):
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    app = win32com.client.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(str(pptx_in), ReadOnly=False, Untitled=False, WithWindow=False)
    try:
        pres.Slides(1).Export(str(png_out), "PNG", 1920, 1080)
    finally:
        pres.Close()


def bake(src_dir):
    src = Path(src_dir)
    tid = src.name
    dst = ASSETS / tid
    dst.mkdir(parents=True, exist_ok=True)
    shutil.copy(src / "template.pptx", dst / "template.pptx")
    if (src / "preview.png").exists():
        shutil.copy(src / "preview.png", dst / "thumb.png")
    desc = json.loads((src / "descriptor.json").read_text(encoding="utf-8"))

    prs = Presentation(str(src / "template.pptx"))
    # blank：清空所有 [key] 文字（保留模板自带固定结构文字），再 COM 渲染
    for sh in iter_shapes(prs.slides[0].shapes):
        if sh.name.startswith("[") and sh.name.endswith("]") and sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""
    tmp = dst / "_blank.pptx"
    prs.save(str(tmp))
    blank_ok = False
    try:
        render_blank_png(tmp, dst / "blank.png")
        blank_ok = True
    except Exception as e:
        print(f"  [{tid}] blank.png COM 渲染失败: {str(e)[:120]}")
    finally:
        try:
            tmp.unlink()
        except Exception:
            pass
    if not (dst / "blank.png").exists() and (dst / "thumb.png").exists():
        shutil.copy(dst / "thumb.png", dst / "blank.png")
        print(f"  [{tid}] blank.png 用 thumb 兜底（预览会含样例文字）")

    smap = {sh.name: sh for sh in iter_shapes(prs.slides[0].shapes)}
    slots = []
    for s in desc["slots"]:
        slot = {k: s.get(k) for k in SLOT_FIELDS}
        sh = smap.get("[" + (s.get("key") or "") + "]")
        # icon 与 image 都不提 style（都走贴图路径，不是文字）
        if s.get("type") not in ("image", "icon") and sh is not None and sh.has_text_frame:
            slot["style"] = extract_style(sh)
        slots.append(slot)

    entry = {
        "id": tid, "name_cn": desc.get("name_cn", tid),
        "when_to_use": desc.get("when_to_use", ""),
        "canvas": desc.get("canvas", [1920, 1080]),
        "thumb": f"ppt-super/assets/{tid}/thumb.png",
        "blank": f"ppt-super/assets/{tid}/blank.png",
        "pptx": f"ppt-super/assets/{tid}/template.pptx",
        "slots": slots,
    }
    tj = SUPER / "templates.json"
    data = json.loads(tj.read_text(encoding="utf-8"))
    data["templates"] = [t for t in data["templates"] if t.get("id") != tid] + [entry]
    tj.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

    n_img = sum(1 for s in slots if s.get("type") == "image")
    n_icon = sum(1 for s in slots if s.get("type") == "icon")
    n_txt = len(slots) - n_img - n_icon
    print(f"  [{tid}] baked: {len(slots)} slots ({n_txt} text + {n_img} image + {n_icon} icon), "
          f"blank={'真实渲染' if blank_ok else '兜底'}")
    return tid


def main():
    srcs = sys.argv[1:] or [DEFAULT_SRC]
    print(f"baking {len(srcs)} template(s) into PPT-SUPER ...")
    for s in srcs:
        bake(s)
    print("DONE")


if __name__ == "__main__":
    main()
