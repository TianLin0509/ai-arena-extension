# enrich_styles.py — 给现有 templates.json 的每个文字 slot 补 style（对齐/字号/加粗/颜色）
# 直接读各模板 template.pptx 提取，不重跑 COM、不动 blank.png，秒级完成。
# 跑法：uv run --with python-pptx python enrich_styles.py
import json, sys
from pathlib import Path
from pptx import Presentation
from build_assets import iter_shapes, extract_style, SUPER
sys.stdout.reconfigure(encoding="utf-8")

ASSETS = SUPER / "assets"
TPL = SUPER / "templates.json"

data = json.loads(TPL.read_text(encoding="utf-8"))
total = 0
for t in data["templates"]:
    prs = Presentation(str(ASSETS / t["id"] / "template.pptx"))
    smap = {sh.name: sh for sh in iter_shapes(prs.slides[0].shapes)}
    n = 0
    for s in t["slots"]:
        if s.get("type") == "image":
            continue
        sh = smap.get("[" + (s.get("key") or "") + "]")
        if sh is None or not sh.has_text_frame:
            continue
        s["style"] = extract_style(sh)
        n += 1
    total += n
    print(f"  [{t['id']}] {n} slots styled", flush=True)

TPL.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
print(f"done: {total} slots across {len(data['templates'])} templates -> {TPL}", flush=True)
